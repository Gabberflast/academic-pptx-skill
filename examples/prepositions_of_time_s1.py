"""
ESL Lesson: Prepositions of Time (IN, ON, AT) — Secondary 1
============================================================
End-to-end example using the ESL Classroom PPTX Skill.
Generates a complete 13-slide PPP lesson from the Presentation.pptx template.

Lesson flow:
  1. Title
  2. Learning Objective
  3. Warm-Up
  4. Grammar Rule: IN
  5. Grammar Rule: ON
  6. Grammar Rule: AT
  7. Summary Comparison Table
  8. Section Divider: Practice
  9. Fill-in-the-Blank (IN / ON / AT)
 10. Multiple Choice (IN / ON / AT)
 11. Error Correction
 12. Sentence Frames: My Daily Routine
 13. Review / Wrap-Up
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, os

# ── Paths ───────────────────────────────────────────────
TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "Presentation.pptx")
OUTPUT   = os.path.join(os.path.dirname(__file__), "Prepositions_of_Time_IN_ON_AT_S1.pptx")

# ── Colours ─────────────────────────────────────────────
C_PRIMARY   = RGBColor(0x0E, 0x28, 0x41)
C_ACCENT    = RGBColor(0x15, 0x60, 0x82)
C_BODY      = RGBColor(0x2D, 0x2D, 0x2D)
C_MUTED     = RGBColor(0x77, 0x77, 0x77)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_RULE      = RGBColor(0xCC, 0xCC, 0xCC)
C_CALLOUT   = RGBColor(0xEB, 0xF3, 0xFA)
C_HIGHLIGHT = RGBColor(0xFF, 0xF2, 0xCC)
C_WORDBANK  = RGBColor(0x5A, 0x40, 0x00)
C_WORDBANK_BORDER = RGBColor(0xE6, 0xC8, 0x00)
C_LIGHT_TXT = RGBColor(0xA0, 0xBB, 0xDD)
C_LIGHTER   = RGBColor(0xCA, 0xDC, 0xFC)
C_FAINT_BG  = RGBColor(0xFA, 0xFA, 0xFA)

# Grammar colours
G_PREP = RGBColor(0xE6, 0x7E, 0x22)  # Orange — prepositions

# Fonts (from template theme)
F_TITLE = "Aptos Display"
F_BODY  = "Aptos"

# ── Helpers ─────────────────────────────────────────────

def add_run(p, text, size=None, font=None, bold=False, italic=False,
            color=None, underline=False):
    run = p.add_run()
    run.text = text
    if size:      run.font.size = size
    if font:      run.font.name = font
    if bold:      run.font.bold = True
    if italic:    run.font.italic = True
    if color:     run.font.color.rgb = color
    if underline: run.font.underline = True
    return run


def add_textbox(slide, left, top, width, height, text="",
                size=Pt(24), font=F_BODY, bold=False, italic=False,
                color=C_BODY, alignment=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    add_run(p, text, size=size, font=font, bold=bold, italic=italic, color=color)
    return tb


def add_rounded_rect(slide, left, top, width, height,
                     fill=C_CALLOUT, border=C_ACCENT, border_pt=1.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def remove_blank_slides(prs):
    """Remove all existing slides from the template."""
    ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    while len(prs.slides._sldIdLst):
        rId = prs.slides._sldIdLst[0].get(f'{ns}id')
        if rId:
            prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


# ── Load template ───────────────────────────────────────
prs = Presentation(TEMPLATE)
remove_blank_slides(prs)

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[0])  # "Title Slide"

title = slide.placeholders[0]
title.text = "Prepositions of Time: IN, ON, AT"
for run in title.text_frame.paragraphs[0].runs:
    run.font.size = Pt(40)
    run.font.bold = True

subtitle = slide.placeholders[1]
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Today: use IN, ON, and AT to talk about when things happen"
p.alignment = PP_ALIGN.CENTER
for run in p.runs:
    run.font.size = Pt(22)

p2 = tf.add_paragraph()
p2.text = "Wednesday, 4 March 2026  ·  S.1  ·  Period 3"
p2.alignment = PP_ALIGN.CENTER
for run in p2.runs:
    run.font.size = Pt(16)
    run.font.color.rgb = C_MUTED


# ============================================================
# SLIDE 2: Learning Objective
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
slide.placeholders[0].text = "Today's Objective"

# Callout box
add_rounded_rect(slide,
    Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.4))
add_textbox(slide,
    Inches(1.8), Inches(2.35), Inches(9.7), Inches(1.1),
    text="Today we will use IN, ON, and AT correctly to talk and write about when things happen.",
    size=Pt(24), font=F_BODY, bold=True, color=C_PRIMARY,
    alignment=PP_ALIGN.CENTER)

# Sub-objectives
obj_box = slide.shapes.add_textbox(
    Inches(1.5), Inches(4.0), Inches(10.3), Inches(2.5))
tf = obj_box.text_frame
tf.word_wrap = True

objectives = [
    ("By the end of this lesson, you can:", True),
    ("• Use IN with months, years, and seasons", False),
    ("• Use ON with days and dates", False),
    ("• Use AT with exact times", False),
    ("• Choose the correct preposition in sentences", False),
]
for i, (text, is_bold) in enumerate(objectives):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    add_run(p, text, size=Pt(24), font=F_BODY, bold=is_bold, color=C_BODY)


# ============================================================
# SLIDE 3: Warm-Up
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[1])  # "Title and Content"
slide.placeholders[0].text = "Warm-Up: When Is It?"

tf = slide.placeholders[1].text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
add_run(p, "Answer these questions with your partner:", size=Pt(24), font=F_BODY, bold=True, color=C_BODY)
p.space_after = Pt(20)

questions = [
    "1. What month is your birthday IN?",
    "2. What day is our English class ON?",
    "3. What time does school start AT?",
]
for q in questions:
    p = tf.add_paragraph()
    p.space_after = Pt(14)
    # Split to colour-code the preposition
    for word in ["IN", "ON", "AT"]:
        if word in q:
            parts = q.split(word)
            add_run(p, parts[0], size=Pt(24), font=F_BODY, color=C_BODY)
            add_run(p, word, size=Pt(24), font=F_BODY, bold=True, color=G_PREP)
            add_run(p, parts[1], size=Pt(24), font=F_BODY, color=C_BODY)
            break

p2 = tf.add_paragraph()
p2.space_after = Pt(0)
p3 = tf.add_paragraph()
add_run(p3, "HINT: ", size=Pt(20), font=F_BODY, bold=True, color=C_ACCENT)
add_run(p3, "The answer uses the preposition at the end of the question!",
        size=Pt(20), font=F_BODY, italic=True, color=C_ACCENT)


# ============================================================
# SLIDE 4: Grammar Rule — IN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
slide.placeholders[0].text = "Prepositions of Time: IN"

# Rule callout box
add_rounded_rect(slide,
    Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.2))

rule_box = slide.shapes.add_textbox(
    Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.0))
tf = rule_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "RULE: ", size=Pt(24), font=F_BODY, bold=True, color=C_PRIMARY)
add_run(p, "Use ", size=Pt(24), font=F_BODY, color=C_PRIMARY)
add_run(p, "IN", size=Pt(26), font=F_BODY, bold=True, color=G_PREP)
add_run(p, " with months, years, seasons, and parts of the day.", size=Pt(24), font=F_BODY, color=C_PRIMARY)

# Examples header
add_textbox(slide,
    Inches(0.92), Inches(3.5), Inches(11.5), Inches(0.4),
    text="EXAMPLES", size=Pt(20), font=F_BODY, bold=True, color=C_ACCENT)

# Examples with colour-coded prepositions
examples_box = slide.shapes.add_textbox(
    Inches(1.3), Inches(4.0), Inches(10.7), Inches(3.0))
tf = examples_box.text_frame
tf.word_wrap = True

examples = [
    ("My birthday is ", "in", " March."),
    ("She was born ", "in", " 2013."),
    ("We go swimming ", "in", " the summer."),
    ("I do my homework ", "in", " the evening."),
]
for i, (before, prep, after) in enumerate(examples):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    add_run(p, before, size=Pt(24), font=F_BODY, color=C_BODY)
    add_run(p, prep, size=Pt(24), font=F_BODY, bold=True, color=G_PREP)
    add_run(p, after, size=Pt(24), font=F_BODY, color=C_BODY)


# ============================================================
# SLIDE 5: Grammar Rule — ON
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.placeholders[0].text = "Prepositions of Time: ON"

add_rounded_rect(slide,
    Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.2))

rule_box = slide.shapes.add_textbox(
    Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.0))
tf = rule_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "RULE: ", size=Pt(24), font=F_BODY, bold=True, color=C_PRIMARY)
add_run(p, "Use ", size=Pt(24), font=F_BODY, color=C_PRIMARY)
add_run(p, "ON", size=Pt(26), font=F_BODY, bold=True, color=G_PREP)
add_run(p, " with days of the week and specific dates.", size=Pt(24), font=F_BODY, color=C_PRIMARY)

add_textbox(slide,
    Inches(0.92), Inches(3.5), Inches(11.5), Inches(0.4),
    text="EXAMPLES", size=Pt(20), font=F_BODY, bold=True, color=C_ACCENT)

examples_box = slide.shapes.add_textbox(
    Inches(1.3), Inches(4.0), Inches(10.7), Inches(3.0))
tf = examples_box.text_frame
tf.word_wrap = True

examples = [
    ("We have English class ", "on", " Monday."),
    ("Her birthday is ", "on", " 15th April."),
    ("I play football ", "on", " Saturdays."),
    ("The test is ", "on", " Friday."),
]
for i, (before, prep, after) in enumerate(examples):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    add_run(p, before, size=Pt(24), font=F_BODY, color=C_BODY)
    add_run(p, prep, size=Pt(24), font=F_BODY, bold=True, color=G_PREP)
    add_run(p, after, size=Pt(24), font=F_BODY, color=C_BODY)


# ============================================================
# SLIDE 6: Grammar Rule — AT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.placeholders[0].text = "Prepositions of Time: AT"

add_rounded_rect(slide,
    Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.2))

rule_box = slide.shapes.add_textbox(
    Inches(1.2), Inches(2.2), Inches(10.9), Inches(1.0))
tf = rule_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "RULE: ", size=Pt(24), font=F_BODY, bold=True, color=C_PRIMARY)
add_run(p, "Use ", size=Pt(24), font=F_BODY, color=C_PRIMARY)
add_run(p, "AT", size=Pt(26), font=F_BODY, bold=True, color=G_PREP)
add_run(p, " with exact times, ", size=Pt(24), font=F_BODY, color=C_PRIMARY)
add_run(p, "night", size=Pt(24), font=F_BODY, bold=True, color=C_PRIMARY)
add_run(p, ", and ", size=Pt(24), font=F_BODY, color=C_PRIMARY)
add_run(p, "the weekend", size=Pt(24), font=F_BODY, bold=True, color=C_PRIMARY)
add_run(p, ".", size=Pt(24), font=F_BODY, color=C_PRIMARY)

add_textbox(slide,
    Inches(0.92), Inches(3.5), Inches(11.5), Inches(0.4),
    text="EXAMPLES", size=Pt(20), font=F_BODY, bold=True, color=C_ACCENT)

examples_box = slide.shapes.add_textbox(
    Inches(1.3), Inches(4.0), Inches(10.7), Inches(3.0))
tf = examples_box.text_frame
tf.word_wrap = True

examples = [
    ("School starts ", "at", " 8:00 a.m."),
    ("I go to bed ", "at", " night."),
    ("We go to the park ", "at", " the weekend."),
    ("The film starts ", "at", " 7:30."),
]
for i, (before, prep, after) in enumerate(examples):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    add_run(p, before, size=Pt(24), font=F_BODY, color=C_BODY)
    add_run(p, prep, size=Pt(24), font=F_BODY, bold=True, color=G_PREP)
    add_run(p, after, size=Pt(24), font=F_BODY, color=C_BODY)


# ============================================================
# SLIDE 7: Summary Comparison Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
slide.placeholders[0].text = "IN, ON, or AT? — Summary"

rows, cols = 5, 4
tbl_shape = slide.shapes.add_table(
    rows, cols, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.5))
table = tbl_shape.table

# Set column widths
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(3.3)
table.columns[2].width = Inches(3.3)
table.columns[3].width = Inches(3.3)

# Header row
headers = ["Preposition", "Use with...", "Signal words", "Example"]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = C_WHITE
            run.font.name = F_BODY
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PRIMARY

# Data
data = [
    ["IN", "months, years,\nseasons, parts of day", "March, 2026,\nsummer, the evening", "I was born in 2013."],
    ["ON", "days, dates", "Monday, 15th April,\nSaturdays", "We play on Fridays."],
    ["AT", "exact times,\nnight, the weekend", "8:00, night,\nthe weekend", "Class starts at 8:00."],
]
# Add an extra row for exceptions
data.append(["REMEMBER", "No preposition with:\nlast, next, this, every", "last Monday,\nnext summer", "I went last Friday.\n(NOT on last Friday)"])

for i, row_data in enumerate(data):
    for j, text in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(17)
                run.font.name = F_BODY
                run.font.color.rgb = C_BODY
                if j == 0:
                    run.font.bold = True
                    run.font.color.rgb = G_PREP if row_data[0] != "REMEMBER" else C_ACCENT


# ============================================================
# SLIDE 8: Section Divider — Practice
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[2])  # "Section Header"
slide.placeholders[0].text = "Time to Practise!"
slide.placeholders[1].text = "Use IN, ON, and AT in exercises."


# ============================================================
# SLIDE 9: Fill-in-the-Blank
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
slide.placeholders[0].text = "Fill in the Blank — IN, ON, or AT?"

# Instruction
add_textbox(slide,
    Inches(0.92), Inches(2.0), Inches(11.5), Inches(0.5),
    text="Complete each sentence with IN, ON, or AT.",
    size=Pt(22), font=F_BODY, italic=True, color=C_BODY)

# Exercise items
items_box = slide.shapes.add_textbox(
    Inches(1.2), Inches(2.7), Inches(11.0), Inches(3.5))
tf = items_box.text_frame
tf.word_wrap = True

items = [
    "1. I wake up __________ 6:30 every morning.",
    "2. My birthday is __________ July.",
    "3. We don't have school __________ Saturday.",
    "4. She reads before bed __________ night.",
    "5. The concert is __________ 21st December.",
    "6. It is very cold __________ winter.",
]
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    add_run(p, item, size=Pt(22), font=F_BODY, color=C_BODY)


# ============================================================
# SLIDE 10: Multiple Choice
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
slide.placeholders[0].text = "Choose the Correct Answer"

# Questions stacked
mc_box = slide.shapes.add_textbox(
    Inches(1.0), Inches(2.1), Inches(11.3), Inches(5.2))
tf = mc_box.text_frame
tf.word_wrap = True

mc_questions = [
    {
        "q": "1. School starts __________ 8 o'clock.",
        "opts": ["A) in", "B) on", "C) at"],
    },
    {
        "q": "2. We have a holiday __________ October.",
        "opts": ["A) in", "B) on", "C) at"],
    },
    {
        "q": "3. The party is __________ Saturday.",
        "opts": ["A) in", "B) on", "C) at"],
    },
    {
        "q": "4. I do exercise __________ the morning.",
        "opts": ["A) in", "B) on", "C) at"],
    },
]

for qi, mc in enumerate(mc_questions):
    # Question
    p = tf.paragraphs[0] if qi == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    add_run(p, mc["q"], size=Pt(22), font=F_BODY, bold=True, color=C_BODY)

    # Options on one line
    p2 = tf.add_paragraph()
    p2.space_after = Pt(14)
    for oi, opt in enumerate(mc["opts"]):
        add_run(p2, opt, size=Pt(22), font=F_BODY, color=C_ACCENT)
        if oi < len(mc["opts"]) - 1:
            add_run(p2, "     ", size=Pt(22), font=F_BODY, color=C_BODY)


# ============================================================
# SLIDE 11: Error Correction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.placeholders[0].text = "Find and Fix the Mistakes"

add_textbox(slide,
    Inches(0.92), Inches(2.0), Inches(11.5), Inches(0.5),
    text="Each sentence has ONE mistake with the preposition. Find it and correct it.",
    size=Pt(22), font=F_BODY, italic=True, color=C_BODY)

items_box = slide.shapes.add_textbox(
    Inches(1.2), Inches(2.7), Inches(11.0), Inches(3.5))
tf = items_box.text_frame
tf.word_wrap = True

errors = [
    "1. I was born at 2013.",
    "2. We have English in Monday.",
    "3. The film starts in 7:30.",
    "4. She goes shopping at Saturdays.",
    "5. It is very hot on summer.",
    "6. I brush my teeth on the morning.",
]
for i, s in enumerate(errors):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    add_run(p, s, size=Pt(22), font=F_BODY, color=C_BODY)

# Hint
add_rounded_rect(slide,
    Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.7),
    fill=C_CALLOUT, border=C_ACCENT)
add_textbox(slide,
    Inches(1.2), Inches(6.3), Inches(10.9), Inches(0.5),
    text="HINT: Think — is it a time, a day, a month, or a season?",
    size=Pt(18), font=F_BODY, italic=True, color=C_ACCENT,
    alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: Sentence Frames — My Daily Routine
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.placeholders[0].text = "My Daily Routine — Use IN, ON, AT"

add_textbox(slide,
    Inches(0.92), Inches(2.0), Inches(11.5), Inches(0.45),
    text="Complete these sentences about YOUR daily routine. Use IN, ON, or AT.",
    size=Pt(22), font=F_BODY, italic=True, color=C_BODY)

frames_box = slide.shapes.add_textbox(
    Inches(1.2), Inches(2.7), Inches(10.9), Inches(2.8))
tf = frames_box.text_frame
tf.word_wrap = True

frames = [
    [("I wake up ", None), ("______", C_ACCENT), (" ____________ every day.", None)],
    [("I have lunch ", None), ("______", C_ACCENT), (" ____________.", None)],
    [("I play with my friends ", None), ("______", C_ACCENT), (" ____________.", None)],
    [("I go to bed ", None), ("______", C_ACCENT), (" ____________.", None)],
    [("I don't have school ", None), ("______", C_ACCENT), (" ____________.", None)],
]
for i, frame in enumerate(frames):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    for text, color in frame:
        add_run(p, text, size=Pt(22), font=F_BODY,
                bold=color is not None, color=color or C_BODY)

# Word bank
add_rounded_rect(slide,
    Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.1),
    fill=C_HIGHLIGHT, border=C_WORDBANK_BORDER)
add_textbox(slide,
    Inches(1.2), Inches(5.75), Inches(10.9), Inches(0.3),
    text="TIME WORDS — choose the correct preposition!", size=Pt(16),
    font=F_BODY, bold=True, color=C_WORDBANK)
add_textbox(slide,
    Inches(1.2), Inches(6.1), Inches(10.9), Inches(0.6),
    text="7:00 a.m.  ·  12:30  ·  Saturdays  ·  the weekend  ·  the afternoon  ·  night  ·  10:00 p.m.  ·  Sunday",
    size=Pt(18), font=F_BODY, color=C_WORDBANK,
    alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13: Review / Wrap-Up
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # "Blank" — custom dark
set_slide_bg(slide, C_PRIMARY)

# Header
add_textbox(slide,
    Inches(0.92), Inches(0.4), Inches(11.5), Inches(0.6),
    text="What We Learned Today", size=Pt(24), font=F_TITLE,
    bold=True, color=C_LIGHT_TXT)

# Accent rule
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(1.05), Inches(11.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = C_ACCENT
line.line.fill.background()

# Key takeaways
takeaways_box = slide.shapes.add_textbox(
    Inches(0.92), Inches(1.4), Inches(11.5), Inches(3.5))
tf = takeaways_box.text_frame
tf.word_wrap = True

takeaways = [
    ("IN", " → months, years, seasons, parts of the day"),
    ("ON", " → days of the week, specific dates"),
    ("AT", " → exact times, night, the weekend"),
    ("NO preposition", " → with last, next, this, every"),
]
for i, (label, detail) in enumerate(takeaways):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(16)
    add_run(p, label, size=Pt(24), font=F_BODY, bold=True, color=G_PREP if label != "NO preposition" else C_LIGHT_TXT)
    add_run(p, detail, size=Pt(24), font=F_BODY, color=C_WHITE)

# Homework
add_textbox(slide,
    Inches(0.92), Inches(5.2), Inches(11.5), Inches(0.4),
    text="HOMEWORK", size=Pt(20), font=F_BODY,
    bold=True, color=C_LIGHT_TXT)

add_textbox(slide,
    Inches(0.92), Inches(5.65), Inches(11.5), Inches(1.2),
    text="Workbook p. 22, exercises 1–3.\nWrite 5 sentences about your weekly routine using IN, ON, and AT.",
    size=Pt(20), font=F_BODY, color=C_LIGHTER)


# ── Save ────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
